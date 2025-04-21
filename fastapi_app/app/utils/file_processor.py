# -*- coding: utf-8 -*-
import os
import base64
import tempfile
import traceback
import io
import logging
import magic
import hashlib
import clamd
from typing import Tuple, Optional, Dict, Any
from pathlib import Path

# File processing dependencies
try:
    import fitz  # type: ignore # PyMuPDF
except ImportError:
    fitz = None

try:
    import ocrmypdf # type: ignore
    # Check if Tesseract is available for ocrmypdf
    # ocrmypdf.check() # Can raise exception if Tesseract is missing
except ImportError:
    ocrmypdf = None
except Exception as ocr_check_e:
     logging.warning(f"ocrmypdf is installed, but there was a problem checking dependencies (e.g., Tesseract): {ocr_check_e}")
     ocrmypdf = None # Treat as unavailable if system dependencies are not met

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None # Handle missing libraries

try:
    import markdown
    from bs4 import BeautifulSoup
except ImportError:
    markdown = None
    BeautifulSoup = None # Handle missing libraries

try:
    import openpyxl  # Excel .xlsx
except ImportError:
    openpyxl = None

try:
    import xlrd  # Excel .xls
except ImportError:
    xlrd = None

try:
    from pptx import Presentation  # PowerPoint
except ImportError:
    Presentation = None

try:
    from striprtf.striprtf import rtf_to_text  # RTF
except ImportError:
    rtf_to_text = None

try:
    import rarfile  # RAR archives
except ImportError:
    rarfile = None

try:
    from odf import text, teletype
    from odf.opendocument import load
except ImportError:
    text = None
    teletype = None
    load = None

try:
    from defusedxml import ElementTree as safe_ET
    import xml.dom.minidom
    from lxml import etree
except ImportError:
    safe_ET = None
    etree = None

try:
    from ebooklib import epub
    from ebooklib.utils import debug
except ImportError:
    epub = None

try:
    import html2text
    from readability import Document
    import html5lib
except ImportError:
    html2text = None
    Document = None

try:
    import xmltodict
    import cssselect
except ImportError:
    xmltodict = None
    cssselect = None

# Logger configuration
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Constants (could be moved to config)
MAX_DOCUMENT_CHARS = 4000 # Max characters returned from a file
MAX_PDF_SIZE_MB = 10  # Max PDF file size in MB
MAX_IMAGE_SIZE_MB = 10 # Max Image file size in MB
MAX_BASE64_SIZE_MB = 15 # Max size for base64 encoded data (~11MB binary)

ALLOWED_MIME_TYPES = [
    'text/plain',
    'application/pdf',
    'image/jpeg',
    'image/png',
    'application/msword',  # .doc
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # .docx
    'application/vnd.ms-excel',  # .xls
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',  # .xlsx
    'text/csv',  # .csv
    'application/vnd.ms-powerpoint',  # .ppt
    'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # .pptx
    'application/rtf',  # .rtf
    'application/zip',  # .zip
    'application/x-rar-compressed',  # .rar
    'application/x-rar',  # alternative MIME type for RAR
    'application/vnd.oasis.opendocument.text',  # .odt
    'application/vnd.oasis.opendocument.spreadsheet',  # .ods
    'application/vnd.oasis.opendocument.presentation',  # .odp
    'text/xml',  # .xml
    'application/xml',  # alternative MIME type for XML
    'text/html',  # .html
    'application/xhtml+xml',  # .xhtml
    'application/epub+zip',  # .epub
]

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB

# --- File Processing Functions ---

def _run_ocr_on_pdf(input_pdf_path: str, output_pdf_path: str) -> bool:
    """Runs OCR on a PDF file using ocrmypdf."""
    if not ocrmypdf:
        logger.warning("ocrmypdf library is not available. Skipping OCR for PDF.")
        return False
    try:
        logger.info(f"Running OCR on PDF file: {input_pdf_path}")
        # Use English and Polish languages, force OCR
        result = ocrmypdf.ocr(input_pdf_path, output_pdf_path, language='eng+pol', force_ocr=True, progress_bar=False)
        if result == 0: # ocrmypdf returns 0 on success
             logger.info(f"OCR completed successfully for: {input_pdf_path}")
             return True
        else:
             logger.error(f"ocrmypdf finished with error code {result} for file {input_pdf_path}")
             return False
    except ocrmypdf.exceptions.TesseractNotFoundError:
         logger.error("Tesseract OCR not found by ocrmypdf.")
         # Change to ValueError to return 400 instead of 500 when OCR is unavailable
         raise ValueError("OCR engine (Tesseract) not found for PDF processing.")
    except Exception as ocr_e:
        logger.error(f"Error during OCR on PDF {input_pdf_path}: {ocr_e}", exc_info=True)
        return False # Return False to attempt extraction without OCR

def safe_extract_text_from_pdf(pdf_data: bytes) -> str:
    """Safely extracts text from PDF data using PyMuPDF, with OCR fallback."""
    if not fitz:
        # fitz is a core dependency, so RuntimeError is appropriate here
        raise RuntimeError("Required library PyMuPDF (fitz) is not installed.")

    text = ""
    temp_input_path = None
    temp_output_path = None

    try:
        # Check file size
        file_size_mb = len(pdf_data) / (1024 * 1024)
        if file_size_mb > MAX_PDF_SIZE_MB:
            raise ValueError(f"PDF file is too large. Maximum size is {MAX_PDF_SIZE_MB} MB.")

        # 1. Try text extraction using PyMuPDF
        try:
            with fitz.open(stream=pdf_data, filetype="pdf") as doc:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    page_text = page.get_text("text") # Extract plain text
                    if page_text:
                        text += page_text + "\n\n"
            logger.info(f"PyMuPDF: Extracted {len(text)} characters of text.")
        except Exception as fitz_e:
             logger.warning(f"Error during text extraction by PyMuPDF: {fitz_e}. Attempting OCR...")
             text = "" # Reset text if PyMuPDF failed

        # 2. If text is empty or insufficient, try OCR
        if not text.strip():
            if ocrmypdf:
                logger.warning("Text extracted by PyMuPDF is empty/insufficient. Attempting OCR with ocrmypdf...")
                try:
                    # Write data to temporary input file
                    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_in:
                        temp_in.write(pdf_data)
                        temp_input_path = temp_in.name

                    # Create path for temporary output file
                    temp_output_fd, temp_output_path = tempfile.mkstemp(suffix=".pdf")
                    os.close(temp_output_fd) # Close file descriptor

                    # Run OCR
                    ocr_success = _run_ocr_on_pdf(temp_input_path, temp_output_path)

                    if ocr_success and os.path.exists(temp_output_path):
                        # Read text from the OCR-processed file using PyMuPDF
                        text = "" # Reset text
                        with fitz.open(temp_output_path) as doc_ocr:
                            for page_num in range(len(doc_ocr)):
                                page = doc_ocr.load_page(page_num)
                                page_text = page.get_text("text")
                                if page_text:
                                    text += page_text + "\n\n"
                        logger.info(f"OCR: Extracted {len(text)} characters after OCR.")
                    else:
                         logger.error(f"OCR failed or output file does not exist: {temp_output_path}")
                         # If OCR fails, still return empty text but log the error
                         text = ""

                except Exception as ocr_pipeline_e:
                     logger.error(f"Error in OCR pipeline for PDF: {ocr_pipeline_e}", exc_info=True)
                     # If OCR fails, still return empty text but log the error
                     text = ""
                finally:
                     # Clean up temporary files
                     if temp_input_path and os.path.exists(temp_input_path):
                         try: os.unlink(temp_input_path)
                         except OSError: logger.warning(f"Could not delete temporary file: {temp_input_path}")
                     if temp_output_path and os.path.exists(temp_output_path):
                         try: os.unlink(temp_output_path)
                         except OSError: logger.warning(f"Could not delete temporary file: {temp_output_path}")
            else: # if ocrmypdf is not available
                 # If ocrmypdf is not available and text is empty, raise an error
                 logger.warning("Text extracted by PyMuPDF is empty, and ocrmypdf is not available for OCR fallback.")
                 raise ValueError("Could not extract text from PDF. OCR library (ocrmypdf) is not available.")

        # Final check if we have any text (after attempting OCR or if OCR wasn't needed/available)
        if not text.strip():
             logger.warning("Failed to extract meaningful text from PDF using PyMuPDF and OCR (if attempted).")
             # Return empty string if nothing found. Error is raised above if OCR was needed but unavailable.

        return text.strip()

    except (ValueError, RuntimeError) as e: # Catch specific errors first
        logger.error(f"Error processing PDF: {e}", exc_info=True)
        # Re-raise as ValueError for the router to handle as 400 or 500 depending on the original type
        raise ValueError(f"Error processing PDF file: {e}") from e
    except Exception as e: # Catch any other unexpected errors
        logger.error(f"Unexpected error processing PDF: {e}", exc_info=True)
        raise RuntimeError(f"Unexpected server error processing PDF: {e}") from e


def safe_extract_text_from_image(image_data: bytes) -> str:
    """Safely extracts text from image data using OCR. Raises ValueError on OCR/dependency errors."""
    if not Image or not pytesseract:
        # Change to ValueError
        raise ValueError("Image processing libraries (Pillow, pytesseract) are not installed.")

    text = ""
    suffix = '.png' # Default
    temp_path = None

    try:
        # Check image format to use correct extension for temporary file
        try:
            img_check = Image.open(io.BytesIO(image_data))
            if img_check.format:
                suffix = f'.{img_check.format.lower()}'
        except Exception:
            logger.warning("Could not determine image format, using default .png")

        # Write data to temporary file
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as temp_file:
            temp_file.write(image_data)
            temp_path = temp_file.name

        # Perform OCR
        image = Image.open(temp_path)
        # Ensure Tesseract path is set (can be done globally)
        # if os.getenv('TESSERACT_CMD'): pytesseract.pytesseract.tesseract_cmd = os.getenv('TESSERACT_CMD')
        text = pytesseract.image_to_string(image, lang='eng+pol') # Added Polish
        logger.info(f"Successfully performed OCR for image (size: {len(image_data)} B).")

    except pytesseract.TesseractNotFoundError:
         logger.error("Tesseract OCR is not installed or not found in PATH.")
         # Change to ValueError
         raise ValueError("OCR engine (Tesseract) is not available for image processing.")
    except Exception as e:
        logger.error(f"Error during image OCR processing: {e}", exc_info=True)
        # Raise as ValueError for consistency
        raise ValueError(f"Error during image OCR processing: {e}") from e
    finally:
        # Delete temporary file
        if temp_path and os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
            except OSError as unlink_e:
                logger.warning(f"Could not delete temporary OCR file: {temp_path}. Error: {unlink_e}")
    return text

def safe_extract_text_from_markdown(md_data: bytes) -> str:
    """Converts Markdown to plain text. Raises ValueError on errors."""
    if not markdown or not BeautifulSoup:
         # Change to ValueError
         raise ValueError("Markdown processing libraries (Markdown, beautifulsoup4) are not installed.")
    try:
        md_text = md_data.decode('utf-8') # Assume UTF-8 for Markdown
        html = markdown.markdown(md_text)
        # Use html.parser as the parser, it's built-in
        soup = BeautifulSoup(html, 'html.parser')
        text = soup.get_text(separator='\n', strip=True)
        return text
    except Exception as e:
        logger.error(f"Error processing Markdown file: {e}", exc_info=True)
        raise ValueError(f"Error processing Markdown file: {e}") from e

def safe_extract_text_from_txt(txt_data: bytes) -> str:
    """Safely reads text from TXT data. Raises ValueError on errors."""
    encodings_to_try = ['utf-8', 'latin-1', 'cp1250', 'cp1252'] # List of encodings to try
    for encoding in encodings_to_try:
        try:
            return txt_data.decode(encoding)
        except UnicodeDecodeError:
            continue # Try next encoding
        except Exception as e:
             logger.error(f"Unexpected error decoding TXT as {encoding}: {e}", exc_info=True)
             raise ValueError(f"Error reading text file: {e}") from e # Raise error if not UnicodeDecodeError

    # If no encoding worked
    logger.error("Could not decode text file using standard encodings.")
    raise ValueError(f"Could not decode text file (tried: {', '.join(encodings_to_try)}).")

def validate_file_type_and_scan(filename: str, content: str) -> Tuple[bool, str, str]:
    """
    Validate file type and scan content with enhanced processing options.
    """
    try:
        decoded_content = base64.b64decode(content)
        mime_type = magic.from_buffer(decoded_content, mime=True)

        if mime_type not in ALLOWED_MIME_TYPES:
            raise ValueError(f"Unsupported file type: {mime_type}")

        # Process content based on MIME type with default options
        if mime_type in ['text/html', 'application/xhtml+xml']:
            processed_content = safe_extract_text_from_html(decoded_content)
        elif mime_type == 'application/epub+zip':
            processed_content = safe_extract_text_from_epub(decoded_content)
        elif mime_type.startswith('text/xml') or mime_type == 'application/xml':
            processed_content = safe_extract_text_from_xml(decoded_content)
        else:
            processed_content = content

        return True, mime_type, processed_content

    except ValueError as e:
        logging.error(f"Validation error: {str(e)}")
        raise
    except Exception as e:
        logging.error(f"Processing error: {str(e)}")
        raise RuntimeError(f"Failed to process file: {str(e)}")

def safe_extract_text_from_excel(excel_data: bytes, mime_type: str) -> str:
    """Safely extracts text from Excel files (.xlsx and .xls)."""
    if mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
        if not openpyxl:
            raise ValueError("openpyxl library is not installed for .xlsx processing")
    elif mime_type == 'application/vnd.ms-excel':
        if not xlrd:
            raise ValueError("xlrd library is not installed for .xls processing")

    text = []
    temp_path = None

    try:
        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx' if 'openxmlformats' in mime_type else '.xls') as temp:
            temp.write(excel_data)
            temp_path = temp.name

        if 'openxmlformats' in mime_type:
            # Process .xlsx
            workbook = openpyxl.load_workbook(temp_path, data_only=True)
            for sheet in workbook.sheetnames:
                ws = workbook[sheet]
                text.append(f"Sheet: {sheet}")
                for row in ws.iter_rows():
                    row_text = ' | '.join(str(cell.value) if cell.value is not None else '' for cell in row)
                    if row_text.strip():
                        text.append(row_text)
        else:
            # Process .xls
            workbook = xlrd.open_workbook(temp_path)
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                text.append(f"Sheet: {sheet.name}")
                for row_idx in range(sheet.nrows):
                    row_text = ' | '.join(str(cell.value) if cell.value else '' for cell in sheet.row(row_idx))
                    if row_text.strip():
                        text.append(row_text)

        return '\n'.join(text)

    except Exception as e:
        logger.error(f"Error processing Excel file: {e}", exc_info=True)
        raise ValueError(f"Error processing Excel file: {e}")
    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
            except OSError:
                pass

def safe_extract_text_from_powerpoint(pptx_data: bytes) -> str:
    """Safely extracts text from PowerPoint files (.pptx)."""
    if not Presentation:
        raise ValueError("python-pptx library is not installed")

    text = []
    temp_path = None

    try:
        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp:
            temp.write(pptx_data)
            temp_path = temp.name

        prs = Presentation(temp_path)
        
        for slide_number, slide in enumerate(prs.slides, 1):
            text.append(f"\nSlide {slide_number}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())

        return '\n'.join(text)

    except Exception as e:
        logger.error(f"Error processing PowerPoint file: {e}", exc_info=True)
        raise ValueError(f"Error processing PowerPoint file: {e}")
    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
            except OSError:
                pass

def safe_extract_text_from_rtf(rtf_data: bytes) -> str:
    """Safely extracts text from RTF files."""
    if not rtf_to_text:
        raise ValueError("striprtf library is not installed")

    try:
        rtf_str = rtf_data.decode('utf-8', errors='ignore')
        text = rtf_to_text(rtf_str)
        return text.strip()
    except Exception as e:
        logger.error(f"Error processing RTF file: {e}", exc_info=True)
        raise ValueError(f"Error processing RTF file: {e}")

def safe_extract_text_from_archive(archive_data: bytes, mime_type: str) -> str:
    """Safely extracts file listing from ZIP/RAR archives."""
    temp_path = None
    try:
        # Save to temporary file
        suffix = '.rar' if 'rar' in mime_type.lower() else '.zip'
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp:
            temp.write(archive_data)
            temp_path = temp.name

        file_list = []
        
        if 'rar' in mime_type.lower():
            if not rarfile:
                raise ValueError("rarfile library is not installed")
            with rarfile.RarFile(temp_path) as rf:
                file_list = rf.namelist()
        else:
            import zipfile
            with zipfile.ZipFile(temp_path) as zf:
                file_list = zf.namelist()

        # Format file listing
        text = "Archive contents:\n" + "\n".join(f"- {name}" for name in file_list)
        return text

    except Exception as e:
        logger.error(f"Error processing archive file: {e}", exc_info=True)
        raise ValueError(f"Error processing archive file: {e}")
    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
            except OSError:
                pass

def safe_extract_text_from_opendocument(odf_data: bytes, mime_type: str) -> str:
    """Safely extracts text from OpenDocument files (ODT, ODS, ODP)."""
    if not all([text, teletype, load]):
        raise ValueError("odfpy library is not installed")

    temp_path = None
    try:
        # Save to temporary file with appropriate extension
        ext_map = {
            'application/vnd.oasis.opendocument.text': '.odt',
            'application/vnd.oasis.opendocument.spreadsheet': '.ods',
            'application/vnd.oasis.opendocument.presentation': '.odp'
        }
        suffix = ext_map.get(mime_type, '.odt')
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp:
            temp.write(odf_data)
            temp_path = temp.name

        doc = load(temp_path)
        
        if mime_type == 'application/vnd.oasis.opendocument.text':
            # Process text document (ODT)
            return teletype.extractText(doc)
            
        elif mime_type == 'application/vnd.oasis.opendocument.spreadsheet':
            # Process spreadsheet (ODS)
            texts = []
            for table in doc.spreadsheet.getElementsByType(text.P):
                if table.firstChild:
                    texts.append(teletype.extractText(table))
            return '\n'.join(texts)
            
        elif mime_type == 'application/vnd.oasis.opendocument.presentation':
            # Process presentation (ODP)
            texts = []
            for slide in doc.presentation.getElementsByType(text.P):
                if slide.firstChild:
                    texts.append(teletype.extractText(slide))
            return '\n\n'.join(texts)
            
        else:
            raise ValueError(f"Unsupported OpenDocument type: {mime_type}")

    except Exception as e:
        logger.error(f"Error processing OpenDocument file: {e}", exc_info=True)
        raise ValueError(f"Error processing OpenDocument file: {e}")
    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
            except OSError:
                pass

def format_extracted_text(text: str, format_type: str = 'default') -> str:
    """
    Formats extracted text according to specified format type.
    
    Args:
        text (str): Text to format
        format_type (str): Type of formatting to apply
            - 'default': Basic cleaning and formatting
            - 'structured': Adds structural elements (headers, sections)
            - 'compact': Removes excessive whitespace
            - 'preserve': Preserves original formatting
            
    Returns:
        str: Formatted text
    """
    if not text:
        return ""
        
    try:
        if format_type == 'structured':
            # Add structural formatting
            lines = text.split('\n')
            formatted_lines = []
            current_section = None
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                    
                # Detect and format headers
                if line.isupper() or len(line) < 50:
                    if current_section:
                        formatted_lines.append('')  # Add spacing between sections
                    formatted_lines.append(f"\n=== {line} ===\n")
                    current_section = line
                else:
                    formatted_lines.append(line)
                    
            text = '\n'.join(formatted_lines)
            
        elif format_type == 'compact':
            # Remove excessive whitespace while preserving paragraph structure
            paragraphs = text.split('\n\n')
            cleaned_paragraphs = []
            
            for para in paragraphs:
                # Clean up whitespace within paragraph
                cleaned = ' '.join(line.strip() for line in para.splitlines() if line.strip())
                if cleaned:
                    cleaned_paragraphs.append(cleaned)
                    
            text = '\n\n'.join(cleaned_paragraphs)
            
        elif format_type == 'preserve':
            # Minimal cleaning while preserving original structure
            return text.rstrip()
            
        else:  # 'default'
            # Basic cleaning
            lines = text.splitlines()
            cleaned_lines = []
            
            for line in lines:
                cleaned = line.strip()
                if cleaned:
                    cleaned_lines.append(cleaned)
                    
            text = '\n'.join(cleaned_lines)
            
        return text.strip()
        
    except Exception as e:
        logger.error(f"Error formatting text: {e}", exc_info=True)
        return text  # Return original text if formatting fails

def safe_extract_text_from_html(content: bytes, formatting_options: Optional[Dict[str, Any]] = None) -> str:
    """
    Extract text from HTML content with customizable formatting options.
    
    Args:
        content: HTML content in bytes
        formatting_options: Dictionary with formatting options:
            - preserve_links: bool - keep links as markdown
            - preserve_images: bool - keep image references
            - preserve_lists: bool - maintain list formatting
            - preserve_tables: bool - maintain table structure
            - extract_article: bool - extract main article content
            - output_format: str - 'text' or 'markdown'
    """
    try:
        if not formatting_options:
            formatting_options = {
                'preserve_links': True,
                'preserve_images': False,
                'preserve_lists': True,
                'preserve_tables': True,
                'extract_article': False,
                'output_format': 'text'
            }

        html_content = content.decode('utf-8')

        if formatting_options.get('extract_article'):
            doc = Document(html_content)
            html_content = doc.summary()

        if formatting_options.get('output_format') == 'markdown':
            h2t = html2text.HTML2Text()
            h2t.ignore_links = not formatting_options.get('preserve_links')
            h2t.ignore_images = not formatting_options.get('preserve_images')
            h2t.ignore_tables = not formatting_options.get('preserve_tables')
            h2t.ignore_emphasis = False
            return h2t.handle(html_content)
        else:
            soup = BeautifulSoup(html_content, 'html5lib')
            if not formatting_options.get('preserve_links'):
                for a in soup.find_all('a'):
                    a.replace_with(a.text)
            if not formatting_options.get('preserve_images'):
                for img in soup.find_all('img'):
                    img.decompose()
            return soup.get_text(separator='\n\n')

    except Exception as e:
        logging.error(f"Error extracting text from HTML: {str(e)}")
        raise RuntimeError(f"Failed to process HTML content: {str(e)}")

def safe_extract_text_from_epub(content: bytes, formatting_options: Optional[Dict[str, Any]] = None) -> str:
    """
    Extract text from EPUB content with formatting options.
    
    Args:
        content: EPUB content in bytes
        formatting_options: Dictionary with formatting options:
            - preserve_chapters: bool - maintain chapter separation
            - preserve_formatting: bool - keep basic formatting
            - extract_metadata: bool - include book metadata
            - include_toc: bool - include table of contents
    """
    try:
        if not formatting_options:
            formatting_options = {
                'preserve_chapters': True,
                'preserve_formatting': True,
                'extract_metadata': True,
                'include_toc': True
            }

        with tempfile.NamedTemporaryFile(delete=False, suffix='.epub') as temp_file:
            temp_file.write(content)
            temp_file.flush()
            
            book = epub.read_epub(temp_file.name)
            
            result = []
            
            if formatting_options.get('extract_metadata'):
                metadata = {
                    'title': book.get_metadata('DC', 'title'),
                    'creator': book.get_metadata('DC', 'creator'),
                    'language': book.get_metadata('DC', 'language')
                }
                result.append(f"Metadata:\n{json.dumps(metadata, indent=2)}\n")

            if formatting_options.get('include_toc'):
                toc = book.toc
                result.append("Table of Contents:")
                for item in toc:
                    result.append(f"- {item.title}")
                result.append("")

            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                if formatting_options.get('preserve_chapters'):
                    soup = BeautifulSoup(item.content, 'html5lib')
                    chapter_title = soup.find(['h1', 'h2', 'h3'])
                    if chapter_title:
                        result.append(f"\n## {chapter_title.text}\n")
                
                if formatting_options.get('preserve_formatting'):
                    h2t = html2text.HTML2Text()
                    h2t.ignore_links = False
                    h2t.ignore_emphasis = False
                    result.append(h2t.handle(item.content.decode('utf-8')))
                else:
                    soup = BeautifulSoup(item.content, 'html5lib')
                    result.append(soup.get_text(separator='\n\n'))

            os.unlink(temp_file.name)
            return '\n'.join(result)

    except Exception as e:
        logging.error(f"Error extracting text from EPUB: {str(e)}")
        raise RuntimeError(f"Failed to process EPUB content: {str(e)}")

def process_xml_with_xslt(xml_content: bytes, xslt_content: bytes) -> str:
    """
    Process XML content using XSLT transformation.
    
    Args:
        xml_content: XML content in bytes
        xslt_content: XSLT stylesheet content in bytes
    """
    try:
        xml_doc = etree.fromstring(xml_content)
        xslt_doc = etree.fromstring(xslt_content)
        transform = XSLT(xslt_doc)
        result = transform(xml_doc)
        return str(result)
    except Exception as e:
        logging.error(f"Error processing XML with XSLT: {str(e)}")
        raise RuntimeError(f"Failed to transform XML with XSLT: {str(e)}")

def safe_extract_text_from_xml(content: bytes, processing_options: Optional[Dict[str, Any]] = None) -> str:
    """
    Extract and format text from XML with advanced processing options.
    
    Args:
        content: XML content in bytes
        processing_options: Dictionary with processing options:
            - use_xpath: str - XPath expression to select nodes
            - format_output: bool - pretty print output
            - convert_to_json: bool - convert to JSON format
            - css_selector: str - use CSS selector instead of XPath
            - xslt_transform: bytes - XSLT stylesheet content
    """
    try:
        if not processing_options:
            processing_options = {
                'format_output': True,
                'convert_to_json': False
            }

        if processing_options.get('xslt_transform'):
            return process_xml_with_xslt(content, processing_options['xslt_transform'])

        parser = etree.XMLParser(resolve_entities=False)
        xml_doc = etree.fromstring(content, parser)

        if processing_options.get('css_selector'):
            translator = cssselect.GenericTranslator()
            xpath_expr = translator.css_to_xpath(processing_options['css_selector'])
            nodes = xml_doc.xpath(xpath_expr)
        elif processing_options.get('use_xpath'):
            nodes = xml_doc.xpath(processing_options['use_xpath'])
        else:
            nodes = [xml_doc]

        if processing_options.get('convert_to_json'):
            if isinstance(nodes, list):
                result = [xmltodict.parse(etree.tostring(node)) for node in nodes]
            else:
                result = xmltodict.parse(etree.tostring(nodes))
            return json.dumps(result, indent=2)
        else:
            if processing_options.get('format_output'):
                return '\n'.join(etree.tostring(node, pretty_print=True, encoding='unicode') for node in nodes)
            else:
                return '\n'.join(etree.tostring(node, encoding='unicode') for node in nodes)

    except Exception as e:
        logging.error(f"Error processing XML content: {str(e)}")
        raise RuntimeError(f"Failed to process XML content: {str(e)}")

def process_uploaded_file_data(file_content_base64: str, filename: str) -> dict:
    """Process uploaded file data and return extracted information"""
    try:
        # Decode base64 content
        file_content = base64.b64decode(file_content_base64)
        
        # Validate file type and scan for viruses
        is_valid, mime_type, _ = validate_file_type_and_scan(filename, file_content_base64)
        if not is_valid:
            raise ValueError(f"Invalid file type: {mime_type}")
            
        # Process file based on type
        extracted_text = ""
        
        if mime_type == 'application/pdf':
            extracted_text = safe_extract_text_from_pdf(file_content)
        elif mime_type.startswith('image/'):
            extracted_text = safe_extract_text_from_image(file_content)
        elif mime_type == 'text/plain':
            extracted_text = safe_extract_text_from_txt(file_content)
        elif mime_type == 'text/markdown':
            extracted_text = safe_extract_text_from_markdown(file_content)
        elif mime_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
            extracted_text = safe_extract_text_from_excel(file_content, mime_type)
        elif mime_type in ['application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
            extracted_text = safe_extract_text_from_powerpoint(file_content)
        elif mime_type == 'application/rtf':
            extracted_text = safe_extract_text_from_rtf(file_content)
        elif mime_type in ['application/zip', 'application/x-rar-compressed', 'application/x-rar']:
            extracted_text = safe_extract_text_from_archive(file_content, mime_type)
        elif mime_type.startswith('application/vnd.oasis.opendocument'):
            extracted_text = safe_extract_text_from_opendocument(file_content, mime_type)
        elif mime_type in ['text/xml', 'application/xml']:
            extracted_text = safe_extract_text_from_xml(file_content)
        elif mime_type in ['text/html', 'application/xhtml+xml']:
            extracted_text = safe_extract_text_from_html(file_content)
        elif mime_type == 'application/epub+zip':
            extracted_text = safe_extract_text_from_epub(file_content)
        else:
            raise ValueError(f"Unsupported file type: {mime_type}")
            
        return {
            "text": extracted_text,
            "mime_type": mime_type,
            "filename": filename
        }
            
    except Exception as e:
        logger.error(f"Error processing file {filename}: {str(e)}")
        raise
